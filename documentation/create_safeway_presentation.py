from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("SafeWay_predefense_presentation.pptx")
ROOT = Path(__file__).resolve().parents[1]


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


BG = RGBColor(13, 22, 30)
PANEL = RGBColor(22, 34, 45)
PANEL_2 = RGBColor(30, 45, 58)
WHITE = RGBColor(246, 250, 252)
MUTED = RGBColor(169, 183, 194)
GREEN = RGBColor(43, 214, 139)
CYAN = RGBColor(58, 184, 255)
AMBER = RGBColor(255, 190, 86)
RED = RGBColor(255, 92, 116)
LINE = RGBColor(68, 91, 106)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def clear(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def add_bg(slide, title=None, section=None, num=None):
    clear(slide)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    # Subtle geometric accents.
    for x, y, w, h, color, trans in [
        (10.7, -0.25, 2.5, 2.5, CYAN, 78),
        (11.55, 5.7, 1.8, 1.8, GREEN, 72),
        (0.55, 6.55, 1.4, 0.24, AMBER, 10),
    ]:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = trans
        shp.line.fill.background()

    if section:
        tx = slide.shapes.add_textbox(Inches(0.72), Inches(0.36), Inches(2.6), Inches(0.3))
        p = tx.text_frame.paragraphs[0]
        p.text = section.upper()
        p.font.name = "Segoe UI Semibold"
        p.font.size = Pt(8.5)
        p.font.color.rgb = GREEN
    if title:
        tx = slide.shapes.add_textbox(Inches(0.72), Inches(0.62), Inches(10.9), Inches(0.68))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI Semibold"
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
    if num:
        tx = slide.shapes.add_textbox(Inches(12.05), Inches(6.92), Inches(0.8), Inches(0.28))
        p = tx.text_frame.paragraphs[0]
        p.text = f"{num:02d}"
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED


def text_box(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Segoe UI Semibold" if bold else "Segoe UI"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tx


def paragraph_box(slide, lines, x, y, w, h, size=15, bullet=True, color=WHITE):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        if bullet:
            p.level = 0
            p.margin_left = Inches(0.2)
    return tx


def panel(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.transparency = 35
    return shp


def stat_card(slide, value, label, x, y, w, h, color=GREEN):
    panel(slide, x, y, w, h, fill=PANEL_2)
    text_box(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.46, size=26, color=color, bold=True)
    text_box(slide, label, x + 0.18, y + 0.75, w - 0.36, h - 0.9, size=10.5, color=MUTED)


def pill(slide, text, x, y, w, color, font_size=10.5):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.fill.transparency = 12
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI Semibold"
    p.font.size = Pt(font_size)
    p.font.color.rgb = BG
    return shp


def link_text(slide, label, url, x, y, w, h, size=16):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.hyperlink.address = url
    run.font.name = "Segoe UI Semibold"
    run.font.size = Pt(size)
    run.font.color.rgb = GREEN
    return tx


def connector(slide, x1, y1, x2, y2, color=GREEN):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(1.6)
    return c


def add_footer(slide):
    text_box(slide, "SafeWay | дипломдық жобаның алдын ала қорғауы | 2026", 0.72, 6.93, 5.2, 0.24, size=8.5, color=MUTED)


def title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    text_box(slide, "ҚАЗАҚСТАН РЕСПУБЛИКАСЫ ҒЫЛЫМ ЖӘНЕ ЖОҒАРЫ БІЛІМ МИНИСТРЛІГІ", 0.7, 0.45, 12.0, 0.3, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(slide, "«КАСПИЙ ҚОҒАМДЫҚ УНИВЕРСИТЕТІ» БМ", 0.7, 0.77, 12.0, 0.3, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "Инженерия институты", 0.7, 1.08, 12.0, 0.28, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(slide, "6В06120 - «Бағдарламалық инженерия» білім беру бағдарламасы", 0.7, 1.36, 12.0, 0.3, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    panel(slide, 1.25, 2.05, 10.85, 2.25, fill=rgb("172632"))
    text_box(slide, "ДИПЛОМДЫҚ ЖОБА", 1.55, 2.28, 10.25, 0.42, size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "«Қауіпсіз маршруттарды құратын\nSafeWay ақылды картасы»", 1.55, 2.88, 10.25, 1.0, size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "Жобаның түрі: мобильді және серверлік бағдарламалық жүйе", 1.55, 3.82, 10.25, 0.32, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    text_box(slide, "Орындаған: ______________________________", 2.15, 4.92, 4.5, 0.32, size=13, color=WHITE)
    text_box(slide, "Жетекші: ________________________________", 6.85, 4.92, 4.5, 0.32, size=13, color=WHITE)
    text_box(slide, "Алматы, 2026", 0.7, 6.6, 12, 0.32, size=13, color=MUTED, align=PP_ALIGN.CENTER)


def slide_goal():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Мақсат, міндеттер және зерттеу негізі", "01 | Жобаның паспорты", 2)
    add_footer(slide)
    panel(slide, 0.75, 1.55, 5.75, 4.95)
    text_box(slide, "Мақсаты", 1.05, 1.82, 2.0, 0.36, size=18, color=GREEN, bold=True)
    paragraph_box(slide, [
        "Ашық картографиялық деректерді, OSRM маршруттауын, жергілікті қауіп қабаттарын және пайдаланушы хабарламаларын біріктіретін SafeWay мобильді жүйесін іске асыру.",
        "Пайдаланушыға тек қысқа жолды емес, қауіп факторлары ескерілген қауіпсіз бағытты ұсыну."
    ], 1.05, 2.34, 5.15, 1.65, size=14.4, bullet=False)
    text_box(slide, "Нысан және пән", 1.05, 4.23, 2.2, 0.34, size=17, color=CYAN, bold=True)
    paragraph_box(slide, [
        "Объект: қала ішінде қауіпсіз бағыт таңдауға көмектесетін мобильді геоақпараттық жүйе.",
        "Пән: маршрут құру, қауіп факторларын есептеу, карта қабаттары, оқиға хабарламалары және баптауларды басқару әдістері."
    ], 1.05, 4.72, 5.15, 1.32, size=13.5, bullet=False, color=MUTED)

    panel(slide, 6.85, 1.55, 5.75, 4.95)
    text_box(slide, "Міндеттер", 7.15, 1.82, 2.0, 0.36, size=18, color=GREEN, bold=True)
    paragraph_box(slide, [
        "Алматы үшін risk zones, safe places және map features деректер моделін құру.",
        "OSRM маршруттарын қауіп және қауіпсіздік қабаттарымен қайта бағалау.",
        "React Native/Expo мобильді интерфейсінде карта, іздеу, қабаттар және маршрутты көрсету.",
        "Пайдаланушы хабарламаларын, аккаунтты және preferences сақтау.",
        "Статистикалық деректер арқылы әлеуметтік өзектілікті негіздеу."
    ], 7.15, 2.35, 5.0, 2.55, size=13.3, bullet=False)
    text_box(slide, "Практикалық құндылығы", 7.15, 5.1, 2.6, 0.32, size=15.5, color=CYAN, bold=True)
    text_box(slide, "Іске қосылатын full-stack прототип: mobile app + Express API + PostgreSQL database.", 7.15, 5.58, 5.0, 0.62, size=13.2, color=MUTED)


def slide_analytics_stats():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Аналитикалық шолу: статистика және әлеуметтік негіз", "02 | Аналитика", 3)
    add_footer(slide)
    text_box(slide, "1234.xlsx деректері Алматы қаласындағы қауіпсіз қозғалыс мәселесінің тұрақты өзектілігін көрсетеді.", 0.78, 1.36, 11.7, 0.48, size=17, color=WHITE)
    stat_card(slide, "7,04", "Алматы көрсеткіші, 2010", 0.82, 2.1, 2.15, 1.3, CYAN)
    stat_card(slide, "5,66", "Алматы көрсеткіші, 2024", 3.22, 2.1, 2.15, 1.3, GREEN)
    stat_card(slide, "7,42", "2010-2024 орташа мән", 5.62, 2.1, 2.15, 1.3, AMBER)
    stat_card(slide, "9,95", "Ең жоғары мән: 2019", 8.02, 2.1, 2.15, 1.3, RED)
    stat_card(slide, "4,36", "Ең төмен мән: 2020", 10.42, 2.1, 2.15, 1.3, GREEN)

    panel(slide, 0.82, 4.0, 5.7, 1.9, fill=rgb("182B36"))
    text_box(slide, "2024 жылғы жоғары тәуекел жас топтары", 1.12, 4.28, 4.8, 0.32, size=16, color=WHITE, bold=True)
    x = 1.15
    for label, col in [("55-59", RED), ("45-49", AMBER), ("20-24", CYAN), ("65-69", GREEN), ("25-29", CYAN)]:
        pill(slide, label, x, 4.92, 0.85, col, 11)
        x += 1.0
    text_box(slide, "Қорытынды: қауіп тек бір аудиторияға тән емес, сондықтан жүйе студенттерге, тұрғындарға, туристерге және кешкі уақытта жүретін қызметкерлерге пайдалы.", 1.12, 5.46, 5.0, 0.28, size=10.6, color=MUTED)

    panel(slide, 7.05, 4.0, 5.25, 1.9, fill=rgb("182B36"))
    text_box(slide, "SafeWay үшін талап", 7.35, 4.28, 2.5, 0.32, size=16, color=WHITE, bold=True)
    paragraph_box(slide, [
        "маршруттың қауіпсіздік индексін көрсету;",
        "жарық, адам көп жүретін дәліз, полиция/аурухана сияқты факторларды картаға шығару;",
        "пайдаланушы хабарламалары арқылы деректі жаңарту."
    ], 7.35, 4.78, 4.55, 0.9, size=11.7, bullet=False, color=MUTED)


def slide_analytics_compare():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Аналитикалық шолу: аналогтар және патенттік контекст", "03 | Аналитика", 4)
    add_footer(slide)
    headers = ["Шешім", "Күшті жағы", "Шектеуі", "SafeWay айырмашылығы"]
    rows = [
        ["Google Maps", "кең карта және навигация", "қауіпсіздік индексі жеке мақсат емес", "risk/safe қабаттармен қайта бағалау"],
        ["2GIS", "қала объектілері жақсы берілген", "route safety логикасы шектеулі", "Алматыға бейімделген қауіп деректері"],
        ["Waze", "қауымдастық хабарламалары", "көбіне автокөлікке бағытталған", "жаяу, bike, scooter профильдері"],
        ["OSM/OSRM", "ашық дерек және routing API", "қауіпсіздік scoring жоқ", "OSRM нәтижесін safety score арқылы сұрыптау"],
    ]
    x0, y0 = 0.78, 1.52
    widths = [1.8, 2.85, 2.85, 4.35]
    row_h = 0.76
    for i, h in enumerate(headers):
        panel(slide, x0 + sum(widths[:i]), y0, widths[i], 0.48, fill=rgb("203340"), radius=False)
        text_box(slide, h, x0 + sum(widths[:i]) + 0.09, y0 + 0.12, widths[i] - 0.18, 0.24, size=10.8, color=GREEN, bold=True)
    for r, row in enumerate(rows):
        for i, cell in enumerate(row):
            fill = rgb("172632") if r % 2 == 0 else rgb("142430")
            panel(slide, x0 + sum(widths[:i]), y0 + 0.48 + r * row_h, widths[i], row_h, fill=fill, radius=False)
            text_box(slide, cell, x0 + sum(widths[:i]) + 0.09, y0 + 0.63 + r * row_h, widths[i] - 0.18, 0.42, size=10.4, color=WHITE if i == 0 else MUTED, bold=(i == 0))

    panel(slide, 0.78, 5.05, 5.55, 1.08, fill=rgb("182B36"))
    text_box(slide, "Ғылыми контекст", 1.05, 5.28, 2.0, 0.3, size=15, color=GREEN, bold=True)
    text_box(slide, "Жоба GIS, spatial analysis, smart map және route scoring ұғымдарына сүйенеді: маршрут геометриясы қауіп/қауіпсіздік қабаттарымен салыстырылады.", 2.75, 5.17, 3.1, 0.58, size=10.8, color=WHITE)
    panel(slide, 6.55, 5.05, 5.88, 1.08, fill=rgb("182B36"))
    text_box(slide, "Патенттік шолу", 6.82, 5.28, 1.75, 0.3, size=15, color=GREEN, bold=True)
    text_box(slide, "Safety metrics, route risk mitigation және safe routing бағытындағы патенттер бар. SafeWay дайын алгоритмді көшірмей, ашық сервистер мен өз scoring формуласын қолданады.", 8.42, 5.17, 3.55, 0.58, size=10.5, color=WHITE)


def slide_problem():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Міндеттің қойылуы: проблема және шешу жолы", "04 | Мәселе", 5)
    add_footer(slide)
    panel(slide, 0.9, 1.55, 3.55, 4.55, fill=rgb("1B2A35"))
    text_box(slide, "Проблема", 1.2, 1.9, 2.4, 0.38, size=21, color=RED, bold=True)
    paragraph_box(slide, [
        "Дәстүрлі навигация көбіне қысқа немесе жылдам жолды таңдайды.",
        "Жаяу жүргінші үшін жарық, адам ағыны, жер асты өткелі, жөндеу, полиция және оқиға хабарламалары маңызды.",
        "Қауіп туралы дерек әртүрлі көзде шашыраңқы орналасқан."
    ], 1.2, 2.55, 2.85, 2.15, size=14, bullet=False)

    connector(slide, 4.65, 3.85, 5.45, 3.85, CYAN)
    panel(slide, 5.7, 1.55, 2.25, 4.55, fill=rgb("142430"))
    text_box(slide, "Әдіс", 6.0, 1.9, 1.5, 0.38, size=21, color=CYAN, bold=True)
    paragraph_box(slide, [
        "1. OSRM нақты жол геометриясын береді.",
        "2. Backend қауіп және қауіпсіз қабаттарды жүктейді.",
        "3. Scoring route geometry мен қауіп нүктелерінің жақындығын есептейді.",
        "4. Route safetyScore бойынша сұрыпталады."
    ], 6.0, 2.55, 1.65, 2.3, size=12.5, bullet=False)

    connector(slide, 8.15, 3.85, 8.95, 3.85, CYAN)
    panel(slide, 9.15, 1.55, 3.25, 4.55, fill=rgb("1B2A35"))
    text_box(slide, "Нәтиже", 9.45, 1.9, 2.4, 0.38, size=21, color=GREEN, bold=True)
    paragraph_box(slide, [
        "Пайдаланушы қауіпсіздік индексі бар ұсынылған маршрутты көреді.",
        "Картада risk zones, user reports, lit streets, crowded corridors және safe places қабаттары қосылады.",
        "Қосымша guest mode және account mode арқылы жұмыс істейді."
    ], 9.45, 2.55, 2.55, 2.1, size=14, bullet=False)
    pill(slide, "Қысқа жол ≠ қауіпсіз жол", 9.52, 5.28, 2.2, AMBER, 10.5)


def slide_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Жүйе архитектурасы", "05 | Шешім", 6)
    add_footer(slide)
    components = [
        ("Mobile app", "React Native / Expo\nкарта, іздеу, профиль,\nқабаттар, маршрут", 0.95, 2.2, CYAN),
        ("Backend API", "Node.js / Express\nвалидация, routing,\nscoring, auth", 4.0, 2.2, GREEN),
        ("PostgreSQL", "risk_zones, safe_places,\nmap_features, reports,\nusers, preferences", 7.05, 2.2, AMBER),
        ("External services", "OSRM, OSM/CARTO,\nNominatim, Overpass", 10.1, 2.2, RED),
    ]
    for title, body, x, y, color in components:
        panel(slide, x, y, 2.35, 1.72, fill=rgb("182B36"))
        pill(slide, title, x + 0.22, y + 0.22, 1.5, color, 9.2)
        text_box(slide, body, x + 0.25, y + 0.72, 1.85, 0.72, size=11.7, color=WHITE)
    connector(slide, 3.3, 3.06, 3.9, 3.06, GREEN)
    connector(slide, 6.35, 3.06, 6.95, 3.06, GREEN)
    connector(slide, 9.4, 3.06, 10.0, 3.06, GREEN)
    panel(slide, 1.35, 4.75, 10.65, 1.0, fill=rgb("142430"))
    text_box(slide, "Деректер ағыны", 1.65, 4.98, 1.7, 0.26, size=14, color=GREEN, bold=True)
    text_box(slide, "Қолданушы әрекеті → HTTP сұраныс → маршрут және қабат деректері → safety scoring → JSON жауап → картадағы визуализация", 3.15, 4.95, 8.25, 0.32, size=13.2, color=WHITE)


def slide_algorithm():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Қауіпсіздік индексін есептеу", "06 | Алгоритм", 7)
    add_footer(slide)
    panel(slide, 0.85, 1.52, 5.15, 4.85, fill=rgb("182B36"))
    text_box(slide, "Scoring логикасы", 1.15, 1.82, 2.3, 0.34, size=18, color=GREEN, bold=True)
    paragraph_box(slide, [
        "riskHits: маршрутқа жақын risk_zones және user_reports.",
        "penalty: severity, қашықтық, avoid таңдауы және түнгі poor_lighting multiplier арқылы өседі.",
        "safeHits: lit_street, crowded_corridor, safe_zone, transport_hub қабаттары.",
        "safetyBoost 16 баллға дейін шектеледі.",
        "Нәтиже 1-100 аралығында safetyScore ретінде беріледі."
    ], 1.15, 2.35, 4.45, 2.7, size=13.2, bullet=False)
    text_box(slide, "score = 100 - penalty - profilePenalty + safetyBoost", 1.15, 5.53, 4.45, 0.28, size=12.4, color=CYAN, bold=True)

    panel(slide, 6.55, 1.52, 5.75, 4.85, fill=rgb("182B36"))
    text_box(slide, "Маршрутты таңдау", 6.85, 1.82, 2.3, 0.34, size=18, color=GREEN, bold=True)
    # Visual score bars.
    labels = [("A: қысқа, бірақ risk көп", 58, RED), ("B: ұзындау, жарық дәліз", 87, GREEN), ("C: орташа балама", 73, CYAN)]
    y = 2.55
    for label, val, col in labels:
        text_box(slide, label, 6.9, y, 2.7, 0.25, size=11.5, color=WHITE)
        panel(slide, 9.7, y + 0.02, 1.72, 0.2, fill=rgb("0F1820"), radius=False)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.7), Inches(y + 0.02), Inches(1.72 * val / 100), Inches(0.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        text_box(slide, str(val), 11.48, y - 0.02, 0.35, 0.25, size=10.2, color=col, bold=True)
        y += 0.7
    text_box(slide, "Егер score айырмасы 6 баллдан көп болса, қауіпсіз маршрут басым. Айырма аз болса, duration бойынша қысқа нұсқа таңдалады.", 6.9, 4.9, 4.75, 0.65, size=12.3, color=MUTED)


def slide_database_api():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Деректер қоры және API", "07 | Backend", 8)
    add_footer(slide)
    tables = [
        ("risk_zones", "қауіп аймақтары\ncategory, severity, radius"),
        ("user_reports", "пайдаланушы хабарламалары\npending/verified status"),
        ("map_features", "жарық көшелер, safe zones,\ntransport hubs"),
        ("safe_places", "полиция, аурухана және\nқоғамдық қауіпсіз орындар"),
        ("users", "аккаунт және JWT auth"),
        ("user_preferences", "JSONB route settings"),
    ]
    positions = [(0.95, 1.65), (3.15, 1.65), (5.35, 1.65), (0.95, 3.55), (3.15, 3.55), (5.35, 3.55)]
    for (name, body), (x, y) in zip(tables, positions):
        panel(slide, x, y, 1.9, 1.25, fill=rgb("182B36"))
        text_box(slide, name, x + 0.15, y + 0.2, 1.55, 0.25, size=12.5, color=GREEN, bold=True)
        text_box(slide, body, x + 0.15, y + 0.58, 1.55, 0.42, size=9.6, color=MUTED)
    panel(slide, 7.85, 1.65, 4.2, 3.15, fill=rgb("142430"))
    text_box(slide, "Негізгі endpoint-тер", 8.15, 1.95, 2.4, 0.32, size=17, color=GREEN, bold=True)
    paragraph_box(slide, [
        "POST /api/routes/safe - recommended және alternatives қайтарады;",
        "GET /api/search?q= - көше/орын іздеу;",
        "GET /api/map/features - карта қабаттары;",
        "POST /api/reports - жаңа оқиға хабарламасы;",
        "POST /api/auth/login және /register - аккаунт."
    ], 8.15, 2.55, 3.55, 1.55, size=11.5, bullet=False)
    text_box(slide, "Zod validation қате координата, профиль және категорияларды API деңгейінде тоқтатады.", 8.15, 4.35, 3.55, 0.3, size=10.5, color=CYAN)


def slide_mobile_demo():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Mobile қосымша және демонстрация", "08 | Demo", 9)
    add_footer(slide)
    # Phone mockup.
    phone = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.34), Inches(3.05), Inches(5.55))
    phone.fill.solid()
    phone.fill.fore_color.rgb = rgb("0A1117")
    phone.line.color.rgb = rgb("526574")
    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.14), Inches(1.62), Inches(2.67), Inches(4.95))
    screen.fill.solid()
    screen.fill.fore_color.rgb = rgb("203340")
    screen.line.fill.background()
    # Map background and route.
    for i, col in enumerate([rgb("2C4857"), rgb("365D67"), rgb("426B74")]):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.22 + i * 0.62), Inches(1.82 + i * 0.5), Inches(2.25 - i * 0.18), Inches(0.08))
        shp.rotation = 18 - i * 22
        shp.fill.solid()
        shp.fill.fore_color.rgb = col
        shp.line.fill.background()
    route_points = [(1.52, 3.12), (1.95, 3.35), (2.25, 3.9), (2.85, 4.28), (2.32, 5.12)]
    for (x1, y1), (x2, y2) in zip(route_points, route_points[1:]):
        seg = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        seg.line.color.rgb = GREEN
        seg.line.width = Pt(4)
    for x, y, col in [(1.45, 2.95, GREEN), (2.75, 4.3, RED), (2.25, 5.15, CYAN)]:
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = col
        marker.line.fill.background()
    panel(slide, 1.35, 5.72, 2.25, 0.62, fill=rgb("101B22"))
    text_box(slide, "Safety score 87", 1.52, 5.9, 1.1, 0.18, size=10.4, color=GREEN, bold=True)
    text_box(slide, "lit street + crowded corridor", 2.45, 5.9, 0.85, 0.18, size=7.3, color=MUTED)

    panel(slide, 4.55, 1.45, 3.45, 4.25, fill=rgb("182B36"))
    text_box(slide, "Көрсетілетін сценарий", 4.85, 1.75, 2.5, 0.34, size=18, color=GREEN, bold=True)
    paragraph_box(slide, [
        "Картаны ашу және Алматы аймағын көрсету.",
        "Орынды іздеу және маршрут нүктесін таңдау.",
        "Профильді таңдау: walk / bike / scooter / drive.",
        "Avoid параметрлері: poor_lighting, underpass, construction.",
        "Қауіпсіз маршрут пен alternatives салыстыру.",
        "Қабаттарды қосу: risk, reports, police, hospitals, lit streets."
    ], 4.85, 2.35, 2.75, 2.4, size=11.8, bullet=False)

    panel(slide, 8.45, 1.45, 3.75, 4.25, fill=rgb("142430"))
    text_box(slide, "Презентациядан ашылатын сілтемелер", 8.75, 1.75, 3.0, 0.34, size=16, color=WHITE, bold=True)
    link_text(slide, "Backend health: http://localhost:4000/health", "http://localhost:4000/health", 8.75, 2.45, 3.1, 0.35, 12.3)
    link_text(slide, "Expo Metro: http://localhost:8081", "http://localhost:8081", 8.75, 2.95, 3.1, 0.35, 12.3)
    link_text(slide, "README.md іске қосу нұсқаулығы", ROOT.joinpath("README.md").as_uri(), 8.75, 3.45, 3.1, 0.35, 12.3)
    text_box(slide, "Демо алдында:", 8.75, 4.18, 1.3, 0.25, size=12, color=AMBER, bold=True)
    text_box(slide, "docker compose up -d postgres\ncd backend && npm run dev\ncd mobile && npm run android", 8.75, 4.55, 2.7, 0.75, size=10.2, color=MUTED)


def slide_testing():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Тестілеу және пайдалану тәртібі", "09 | Нәтиже", 10)
    add_footer(slide)
    steps = [
        ("1", "PostgreSQL", "docker compose up -d postgres\nnpm run db:setup", GREEN),
        ("2", "Backend", "npm run dev\nGET /health", CYAN),
        ("3", "Mobile", "npm start немесе npm run android\nExpo Go / Android device", AMBER),
        ("4", "Қолданушы сценарийі", "іздеу → route → layers → report → preferences", RED),
    ]
    x = 0.85
    for n, title, body, col in steps:
        panel(slide, x, 1.75, 2.75, 3.9, fill=rgb("182B36"))
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.25), Inches(2.05), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = col
        circ.line.fill.background()
        text_box(slide, n, x + 0.25, 2.15, 0.55, 0.18, size=14, color=BG, bold=True, align=PP_ALIGN.CENTER)
        text_box(slide, title, x + 0.95, 2.08, 1.45, 0.32, size=15, color=WHITE, bold=True)
        text_box(slide, body, x + 0.28, 3.0, 2.15, 1.05, size=11.4, color=MUTED)
        x += 3.0
    text_box(slide, "Тексеру нәтижесі: жүйе локалды ортада backend, database және mobile клиент ретінде толық байланысады; қауіпсіз маршрут API арқылы есептеліп, картада визуалды көрсетіледі.", 0.95, 6.05, 11.2, 0.45, size=14.3, color=WHITE)


def slide_value():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Практикалық нәтиже және апробация", "10 | Нәтиже", 11)
    add_footer(slide)
    panel(slide, 0.9, 1.55, 5.4, 4.9, fill=rgb("182B36"))
    text_box(slide, "Дайын нәтиже", 1.2, 1.85, 2.2, 0.35, size=19, color=GREEN, bold=True)
    paragraph_box(slide, [
        "React Native / Expo мобильді қосымшасы;",
        "Express.js backend API;",
        "PostgreSQL migrations және demo seed деректері;",
        "JWT арқылы optional account mode;",
        "Guest preferences және user reports сақтау;",
        "README арқылы қайта іске қосу тәртібі."
    ], 1.2, 2.45, 4.6, 2.3, size=13.5, bullet=False)
    text_box(slide, "Апробация: локалды ортада Android/Expo сценарийі арқылы маршрут құру, қабаттарды көрсету және API health check тексерілді.", 1.2, 5.3, 4.55, 0.62, size=12.4, color=CYAN)

    panel(slide, 6.9, 1.55, 5.4, 4.9, fill=rgb("142430"))
    text_box(slide, "Кеңейту бағыты", 7.2, 1.85, 2.2, 0.35, size=19, color=GREEN, bold=True)
    paragraph_box(slide, [
        "PostGIS арқылы spatial index және ST_DWithin енгізу.",
        "Verified data: әкімдік, полиция, жарықтандыру қызметі.",
        "Admin dashboard және user reports модерациясы.",
        "Offline cache және соңғы маршрутты сақтау.",
        "Жеке OSRM/Nominatim instance және production HTTPS.",
        "Push notification және аналитикалық dashboard."
    ], 7.2, 2.45, 4.55, 2.5, size=13.2, bullet=False)


def slide_conclusion():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Қорытынды", "11 | Финал", 12)
    add_footer(slide)
    panel(slide, 1.1, 1.6, 11.0, 4.75, fill=rgb("182B36"))
    paragraph_box(slide, [
        "SafeWay қалалық қауіпсіздік пен мобильді навигацияны біріктіретін қолданбалы бағдарламалық шешім ретінде әзірленді.",
        "Жоба ең қысқа маршрутты ғана емес, қауіп және қауіпсіздік факторларымен қайта бағаланған route recommendation ұсынады.",
        "Негізгі нәтиже - React Native mobile app, Express backend және PostgreSQL деректер қоры бар жұмыс істейтін full-stack прототип.",
        "Статистикалық деректер тақырыптың әлеуметтік өзектілігін растады, ал жүйе әрі қарай қалалық пилот немесе кампус қауіпсіз навигациясы ретінде кеңейтіле алады."
    ], 1.55, 2.1, 10.0, 2.45, size=17, bullet=False)
    text_box(slide, "Мақсат орындалды: қауіпсіз маршрут құру, карта қабаттарын көрсету, хабарлама сақтау, іздеу және мобильді навигация элементтері бір жүйеге біріктірілді.", 1.55, 5.25, 10.0, 0.44, size=15, color=GREEN, bold=True)


def slide_sources():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "Пайдаланылған дереккөздер", "12 | Sources", 13)
    add_footer(slide)
    sources = [
        "1234.xlsx. ҚР жаяу жүргіншілерге қатысты статистикасы, 2010-2024.",
        "Express.js Routing Guide - https://expressjs.com/en/guide/routing.html",
        "PostgreSQL JSON Types - https://www.postgresql.org/docs/current/datatype-json.html",
        "Expo SDK 54 және Location API - https://docs.expo.dev/",
        "React Native Maps - https://github.com/react-native-maps/react-native-maps",
        "OSRM API Documentation - https://project-osrm.org/docs/v5.6.4/api/",
        "OpenStreetMap / Nominatim / Overpass usage policies.",
        "CARTO Basemaps - https://docs.carto.com/carto-user-manual/maps/basemaps",
        "Google Patents: US8315792B2, US9932033B2, US10563994B2, EP2372305A2.",
        "SafeWay бастапқы коды: backend, mobile, db migrations және README.md."
    ]
    paragraph_box(slide, sources, 1.0, 1.45, 11.2, 4.8, size=12.3, bullet=False, color=WHITE)
    text_box(slide, "Назарларыңызға рақмет", 1.0, 6.28, 11.2, 0.36, size=21, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


for builder in [
    title_slide,
    slide_goal,
    slide_analytics_stats,
    slide_analytics_compare,
    slide_problem,
    slide_architecture,
    slide_algorithm,
    slide_database_api,
    slide_mobile_demo,
    slide_testing,
    slide_value,
    slide_conclusion,
    slide_sources,
]:
    builder()


prs.save(OUT)
print(OUT)
